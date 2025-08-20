# parsers/pptx_parser.py
from typing import List, Optional
from models import DocElement
from pathlib import Path
from pptx import Presentation

def parse_pptx(path: str, category: Optional[str]=None) -> List[DocElement]:
    prs = Presentation(path)
    stem = Path(path).stem
    out: List[DocElement] = []
    doc_id = stem
    order = 0

    for sidx, slide in enumerate(prs.slides, start=1):
        # Заголовки и текстовые фреймы
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                txt = shape.text.strip()
                if txt:
                    out.append(DocElement(
                        category=category, doc_path=path, doc_id=doc_id, source_type="pptx",
                        element_type="paragraph", text=txt, order=order,
                        metadata={"slide": sidx}
                    ))
                    order += 1
            if hasattr(shape, "has_table") and shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.replace("\n"," ").strip() for c in row.cells]
                    out.append(DocElement(
                        category=category, doc_path=path, doc_id=doc_id, source_type="pptx",
                        element_type="table", text="\t".join(cells), order=order,
                        metadata={"slide": sidx}
                    ))
                    order += 1
        # Отметим границу слайда (метаданные)
        out.append(DocElement(
            category=category, doc_path=path, doc_id=doc_id, source_type="pptx",
            element_type="slide", text=f"--- Слайд {sidx} ---", order=order
        ))
        order += 1

    return out